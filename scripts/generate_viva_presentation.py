from __future__ import annotations

import subprocess
from pathlib import Path
from textwrap import wrap

from PIL import Image, ImageOps, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ROOT = Path('/Users/nathanbrown-bennett/Abokar/mentra')
MEDIA = ROOT / 'docs' / 'media'
OUTPUT_DIR = ROOT / 'docs' / 'viva'
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

PRESENTATION_PATH = OUTPUT_DIR / 'mentra-viva-presentation.pptx'
VIDEO_WEBM = MEDIA / '16-app-demo-refresh-2026.webm'
VIDEO_MP4 = MEDIA / '16-app-demo-refresh-2026.mp4'
VIDEO_POSTER = MEDIA / '15-learning-feedback-transparency-2026.png'

TITLE_BLUE = RGBColor(15, 23, 42)
ACCENT = RGBColor(79, 70, 229)
ACCENT_2 = RGBColor(14, 165, 233)
SOFT_BG = RGBColor(248, 250, 252)
TEXT = RGBColor(30, 41, 59)
MUTED = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(22, 163, 74)
ORANGE = RGBColor(249, 115, 22)
RED = RGBColor(220, 38, 38)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide, color=SOFT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_bar(slide, title, subtitle=None, dark=False):
    bar_color = TITLE_BLUE if dark else ACCENT
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.62))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bar_color
    shape.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.45), Inches(0.12), Inches(10.5), Inches(0.28))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = WHITE

    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(8.1), Inches(0.12), Inches(4.7), Inches(0.28))
        tf2 = tx2.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(10)
        r2.font.color.rgb = WHITE


def add_footer(slide, text='Mentra viva presentation'):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(7.0), Inches(12.2), Inches(0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED


def add_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, TITLE_BLUE)

    accent = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.6), Inches(3.0), Inches(0.55))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_2
    accent.line.fill.background()
    tf = accent.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = 'Dissertation Viva'
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    title = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(6.7), Inches(2.0))
    tf = title.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = 'Mentra'
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = 'AI-driven personal tutor with evidence, explainability, guardrails, and transparent feedback'
    r2.font.size = Pt(18)
    r2.font.color.rgb = RGBColor(226, 232, 240)

    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = 'React + FastAPI + SQLAlchemy + local Ollama-backed feedback'
    r3.font.size = Pt(14)
    r3.font.color.rgb = RGBColor(191, 219, 254)

    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(4.15), Inches(2.1), Inches(0.5))
    chip.fill.solid()
    chip.fill.fore_color.rgb = RGBColor(30, 41, 59)
    chip.line.fill.background()
    chip_tf = chip.text_frame
    chip_tf.clear()
    p = chip_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'Local demo build'
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = WHITE

    for idx, (x, y, w, h, c) in enumerate([
        (8.3, 0.8, 4.0, 2.1, ACCENT_2),
        (8.8, 2.3, 3.7, 1.6, ACCENT),
        (7.8, 3.55, 4.6, 2.7, RGBColor(30, 41, 59)),
    ]):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = c
        shape.line.fill.background()

    img = slide.shapes.add_picture(str(MEDIA / '12-dashboard-evidence-2026.png'), Inches(8.05), Inches(3.72), width=Inches(4.25))
    img.line.color.rgb = WHITE
    add_footer(slide, 'Mentra - viva presentation')


def add_bullets_slide(title, subtitle, bullets, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_bar(slide, title, subtitle)

    left = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(6.0), Inches(5.6))
    tf = left.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(12)
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT

    if notes:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.0), Inches(5.4), Inches(5.7))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(226, 232, 240)
        tx = slide.shapes.add_textbox(Inches(7.55), Inches(1.25), Inches(4.7), Inches(5.0))
        tf2 = tx.text_frame
        tf2.word_wrap = True
        tf2.clear()
        for idx, line in enumerate(notes):
            p = tf2.paragraphs[0] if idx == 0 else tf2.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(14)
            p.font.color.rgb = MUTED
            p.space_after = Pt(8)

    add_footer(slide)


def add_image_slide(title, subtitle, image_path, caption, extra_text=None, zoom=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_bar(slide, title, subtitle)

    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.0), Inches(8.1), Inches(5.75))
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = RGBColor(226, 232, 240)

    if zoom:
        slide.shapes.add_picture(str(image_path), Inches(0.8), Inches(1.25), width=Inches(7.5))
    else:
        slide.shapes.add_picture(str(image_path), Inches(0.8), Inches(1.2), width=Inches(7.5))

    text_box = slide.shapes.add_textbox(Inches(8.85), Inches(1.2), Inches(3.7), Inches(4.9))
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = caption
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT
    p.space_after = Pt(12)

    if extra_text:
        for item in extra_text:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = MUTED
            p.space_after = Pt(7)

    add_footer(slide)


def add_architecture_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_bar(slide, 'System Architecture', 'How the app is built')

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.0), Inches(12.2), Inches(5.9))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(226, 232, 240)

    slide.shapes.add_picture(str(MEDIA / 'fyp-architecture-diagram.png'), Inches(0.8), Inches(1.15), width=Inches(5.6))
    slide.shapes.add_picture(str(MEDIA / 'fyp-use-case-diagram.png'), Inches(6.65), Inches(1.15), width=Inches(5.9))

    tx = slide.shapes.add_textbox(Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.45))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'React frontend | FastAPI backend | SQLAlchemy persistence | local Ollama feedback service | seeded demo data'
    r.font.size = Pt(13)
    r.font.color.rgb = MUTED

    add_footer(slide)


def add_flow_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_bar(slide, 'Feedback Flow', 'Question -> assessment -> explainable output')

    slide.shapes.add_picture(str(MEDIA / 'fyp-sequence-diagram.png'), Inches(0.75), Inches(1.05), width=Inches(6.0))

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.15), Inches(5.45), Inches(5.55))
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = RGBColor(226, 232, 240)

    tx = slide.shapes.add_textbox(Inches(7.55), Inches(1.35), Inches(4.9), Inches(5.0))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.clear()
    lines = [
        'The learner answers a topic question.',
        'The backend evaluates correctness and stores the answer.',
        'LLM feedback is generated locally through Ollama.',
        'The response includes explanation, key concepts, confidence, and next step.',
        'Progress and evidence endpoints update the dissertation metrics layer.',
    ]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(17)
        p.font.color.rgb = TEXT
        p.space_after = Pt(12)
    add_footer(slide)


def add_video_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_bar(slide, 'Live Demo Video', 'Embedded app walkthrough')

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.0), Inches(12.2), Inches(5.9))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(226, 232, 240)

    slide.shapes.add_textbox(Inches(0.9), Inches(1.25), Inches(11.2), Inches(0.4))
    note = slide.shapes.add_textbox(Inches(0.9), Inches(1.15), Inches(11.4), Inches(0.5))
    tf = note.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'The embedded demo shows login, dashboard evidence, progress analytics, and the learning feedback transparency flow.'
    r.font.size = Pt(14)
    r.font.color.rgb = MUTED

    slide.shapes.add_movie(
        str(VIDEO_MP4),
        Inches(1.1), Inches(1.7), Inches(11.1), Inches(4.6),
        poster_frame_image=str(VIDEO_POSTER),
        mime_type='video/mp4',
    )

    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.25), Inches(6.45), Inches(2.0), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT_2
    badge.line.fill.background()
    tfb = badge.text_frame
    tfb.clear()
    p = tfb.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'Demo video'
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = WHITE

    add_footer(slide)


def add_closing_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, TITLE_BLUE)

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(5.8), Inches(4.5))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = 'Key Takeaways'
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = WHITE

    for line in [
        'Mentra demonstrates adaptive tutoring with explainable feedback.',
        'The dashboard and progress views surface measurable learning evidence.',
        'Guardrails and transparency screens show the system is designed for responsible use.',
        'The local Ollama-backed runtime keeps the demo self-contained for viva delivery.',
    ]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(17)
        p.font.color.rgb = RGBColor(226, 232, 240)
        p.space_after = Pt(10)

    slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.0), Inches(5.5), Inches(4.7))
    slide.shapes.add_picture(str(MEDIA / '14-learning-guardrail-2026.png'), Inches(7.2), Inches(1.25), width=Inches(5.0))
    slide.shapes.add_textbox(Inches(7.1), Inches(6.1), Inches(5.4), Inches(0.4))
    tx = slide.shapes.add_textbox(Inches(7.1), Inches(6.15), Inches(5.4), Inches(0.4))
    tfx = tx.text_frame
    tfx.clear()
    p = tfx.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'Questions and discussion'
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = WHITE


def ensure_video_mp4():
    if VIDEO_MP4.exists() and VIDEO_MP4.stat().st_size > 0:
        return
    subprocess.run([
        'ffmpeg', '-y', '-i', str(VIDEO_WEBM),
        '-c:v', 'libx264', '-pix_fmt', 'yuv420p',
        '-c:a', 'aac', '-movflags', '+faststart',
        str(VIDEO_MP4),
    ], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)


def main():
    ensure_video_mp4()

    add_title_slide()
    add_bullets_slide(
        'Why Mentra Exists',
        'Problem framing and contribution',
        [
            'Traditional e-learning tools give limited feedback after an answer is submitted.',
            'Mentra adds adaptive questioning, immediate explanations, and progress evidence.',
            'The system is designed for a viva-ready dissertation demonstration, not just a static prototype.',
        ],
        notes=[
            'Use this slide to frame the problem space.',
            'The emphasis is on personalisation, explanation, and measurable learning impact.',
            'Mentra is a self-contained local demo with seeded data and a local LLM runtime.',
        ],
    )
    add_architecture_slide()
    add_flow_slide()
    add_image_slide(
        'Evidence on the Dashboard',
        'Screenshot: dashboard evidence + explainability',
        MEDIA / '12-dashboard-evidence-2026.png',
        'Dashboard surfaces measurable outcomes and recommendation explainability.',
        [
            'Baseline vs current accuracy are exposed on the landing dashboard.',
            'Reason codes and confidence intervals make recommendation logic legible.',
        ],
    )
    add_image_slide(
        'Progress as Evidence',
        'Screenshot: progress analytics and experiment events',
        MEDIA / '13-progress-evidence-2026.png',
        'Progress page shows learning evidence, topic status, and experiment telemetry.',
        [
            'The evidence layer shows baseline, current accuracy, and confidence trend.',
            'Topic rows summarise attempts, correctness, and mastery status.',
        ],
    )
    add_image_slide(
        'Guardrails and Transparency',
        'Screenshots: safe learning state and explanation rubric',
        MEDIA / '14-learning-guardrail-2026.png',
        'Guardrail state prevents misuse and sets expectations before feedback is generated.',
        [
            'The learning page includes explicit ethical guardrails.',
            'The transparency flow shows rubric-backed feedback after submission.',
        ],
    )
    add_video_slide()
    add_closing_slide()

    prs.save(PRESENTATION_PATH)
    print(f'Saved presentation: {PRESENTATION_PATH}')


if __name__ == '__main__':
    main()
